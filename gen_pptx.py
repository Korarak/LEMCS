#!/usr/bin/env python3
"""Generate LEMCS Executive + Technical Presentation (PPTX) — Content Style C"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── Palette ───────────────────────────────────────────────────────────────────
NAVY    = RGBColor(0x0B, 0x15, 0x26)
NAVY2   = RGBColor(0x08, 0x14, 0x22)
CARD    = RGBColor(0x12, 0x20, 0x36)
CARD2   = RGBColor(0x0F, 0x1C, 0x30)
TEAL    = RGBColor(0x00, 0xC9, 0xB1)
TEAL2   = RGBColor(0x00, 0xE5, 0xC8)
BLUE    = RGBColor(0x1E, 0x90, 0xFF)
GREEN   = RGBColor(0x2E, 0xCC, 0x71)
YELLOW  = RGBColor(0xF0, 0xC0, 0x40)
RED     = RGBColor(0xE7, 0x4C, 0x3C)
PURPLE  = RGBColor(0x9B, 0x59, 0xB6)
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT   = RGBColor(0xE8, 0xF0, 0xFE)
MUTED   = RGBColor(0x8F, 0xA3, 0xBC)
DKBLUE  = RGBColor(0x08, 0x1E, 0x38)
DKROW   = RGBColor(0x10, 0x1C, 0x30)

# ── Slide dimensions (16:9 widescreen) ───────────────────────────────────────
W = Inches(13.333)
H = Inches(7.5)
M = Inches(0.6)

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H
BLANK = prs.slide_layouts[6]

TOTAL_SLIDES = 22

# ── Primitives ────────────────────────────────────────────────────────────────

def slide(bg=NAVY):
    sl = prs.slides.add_slide(BLANK)
    f = sl.background.fill
    f.solid()
    f.fore_color.rgb = bg
    return sl

def box(sl, x, y, w, h, fill, line_color=None, line_pt=0):
    sh = sl.shapes.add_shape(1, x, y, w, h)
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill
    if line_color:
        sh.line.color.rgb = line_color
        sh.line.width = Pt(line_pt or 1)
    else:
        sh.line.fill.background()
    return sh

def txt(sl, text, x, y, w, h, size=13, color=LIGHT, bold=False,
        align=PP_ALIGN.LEFT, italic=False):
    tb = sl.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    p  = tf.paragraphs[0]
    p.alignment = align
    r  = p.add_run()
    r.text        = text
    r.font.size   = Pt(size)
    r.font.color.rgb = color
    r.font.bold   = bold
    r.font.italic = italic
    r.font.name   = 'Sarabun'
    return tb

def mtxt(sl, lines, x, y, w, h, def_size=12, def_color=LIGHT, def_align=PP_ALIGN.LEFT):
    """Multi-paragraph textbox. Each item = str or dict."""
    tb = sl.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    first = True
    for item in lines:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        if isinstance(item, str):
            p.alignment = def_align
            r = p.add_run()
            r.text = item
            r.font.size = Pt(def_size)
            r.font.color.rgb = def_color
            r.font.name = 'Sarabun'
        else:
            p.alignment = item.get('align', def_align)
            if item.get('sb'): p.space_before = Pt(item['sb'])
            if item.get('sa'): p.space_after  = Pt(item['sa'])
            r = p.add_run()
            r.text = item.get('t', '')
            r.font.size  = Pt(item.get('s', def_size))
            r.font.color.rgb = item.get('c', def_color)
            r.font.bold  = item.get('b', False)
            r.font.italic = item.get('i', False)
            r.font.name  = 'Sarabun'
    return tb

def header(sl, title, label='', n=None):
    box(sl, 0, 0, W, Inches(0.07), TEAL)
    txt(sl, title, M, Inches(0.15), Inches(10.5), Inches(0.72),
        size=26, color=WHITE, bold=True)
    box(sl, M, Inches(0.9), Inches(0.9), Inches(0.04), TEAL)
    if label:
        txt(sl, label, Inches(10.5), Inches(0.2), Inches(2.6), Inches(0.4),
            size=9, color=MUTED, align=PP_ALIGN.RIGHT)
    if n:
        txt(sl, f'{n} / {TOTAL_SLIDES}',
            W - Inches(1.6), H - Inches(0.42), Inches(1.3), Inches(0.32),
            size=9, color=MUTED, align=PP_ALIGN.RIGHT)

def sec_slide(num_str, title, subtitle, accent=TEAL, n=None, bg=RGBColor(0x07, 0x14, 0x24)):
    sl = slide(bg)
    box(sl, 0, 0, Inches(0.55), H, accent)
    txt(sl, num_str, Inches(1.2), Inches(2.3), Inches(10), Inches(0.45),
        size=11, color=accent, bold=True)
    txt(sl, title, Inches(1.2), Inches(2.82), Inches(11), Inches(1.15),
        size=44, color=WHITE, bold=True)
    txt(sl, subtitle, Inches(1.2), Inches(4.0), Inches(11), Inches(0.45),
        size=13, color=MUTED)
    if n:
        txt(sl, f'{n} / {TOTAL_SLIDES}',
            W - Inches(1.6), H - Inches(0.42), Inches(1.3), Inches(0.32),
            size=9, color=MUTED, align=PP_ALIGN.RIGHT)
    return sl

def card(sl, x, y, w, h, accent=TEAL, title='', body='', icon=''):
    box(sl, x, y, w, h, CARD)
    box(sl, x, y, Inches(0.055), h, accent)
    cx = x + Inches(0.18)
    cy = y + Inches(0.13)
    cw = w - Inches(0.25)
    row_h = Inches(0.35)
    if icon:
        txt(sl, icon, cx, cy, Inches(0.45), row_h, size=16)
        cx2 = cx + Inches(0.46)
    else:
        cx2 = cx
    if title:
        txt(sl, title, cx2, cy, cw - Inches(0.46) if icon else cw,
            row_h, size=11, color=accent, bold=True)
    if body:
        by = y + Inches(0.13) + row_h + Inches(0.06)
        bh = h - row_h - Inches(0.28)
        mtxt(sl, [{'t': line, 's': 9, 'c': MUTED} for line in body.split('\n')],
             cx, by, w - Inches(0.28), bh)

def stat(sl, x, y, w, h, number, label, accent=TEAL2):
    box(sl, x, y, w, h, CARD)
    box(sl, x, y, w, Inches(0.05), accent)
    txt(sl, number, x, y + Inches(0.08), w, Inches(0.58),
        size=26, color=accent, bold=True, align=PP_ALIGN.CENTER)
    txt(sl, label,  x, y + Inches(0.62), w, Inches(0.42),
        size=9, color=MUTED, align=PP_ALIGN.CENTER)

def tbl(sl, x, y, w, h, rows, col_widths, hdr_color=DKBLUE):
    t = sl.shapes.add_table(len(rows), len(rows[0]), x, y, w, h).table
    for ci, cw in enumerate(col_widths):
        t.columns[ci].width = cw
    for ri, row_data in enumerate(rows):
        for ci, cell_text in enumerate(row_data):
            cell = t.cell(ri, ci)
            cell.text = str(cell_text)
            p = cell.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.LEFT
            runs = p.runs
            r = runs[0] if runs else p.add_run()
            r.font.name = 'Sarabun'
            if ri == 0:
                r.font.size  = Pt(11)
                r.font.bold  = True
                r.font.color.rgb = TEAL2
                cell.fill.solid()
                cell.fill.fore_color.rgb = hdr_color
            else:
                r.font.size  = Pt(10)
                r.font.color.rgb = LIGHT
                r.font.bold  = (ci == 0)
                cell.fill.solid()
                cell.fill.fore_color.rgb = CARD if ri % 2 == 0 else DKROW
    return t

# ═══════════════════════════════════════════════════════════════════════════════
# S01 — Cover
# ═══════════════════════════════════════════════════════════════════════════════
sl = slide(NAVY2)
box(sl, 0,            0, Inches(0.55), H, TEAL)
box(sl, Inches(0.55), 0, Inches(0.28), H, RGBColor(0x00, 0x78, 0x6E))

# Badge ribbon
box(sl, Inches(1.2), Inches(0.62), Inches(11.5), Inches(0.42), DKBLUE)
txt(sl, 'LOEI EDUCATIONAL MINDCARE SYSTEM  ·  ระบบดิจิทัลประเมินสุขภาพจิตนักเรียน จังหวัดเลย',
    Inches(1.4), Inches(0.68), Inches(11.2), Inches(0.32),
    size=9, color=TEAL, bold=True)

txt(sl, 'LEMCS', Inches(1.2), Inches(1.3), Inches(11), Inches(1.9),
    size=100, color=TEAL2, bold=True)
txt(sl, 'ระบบดิจิทัลประเมินสุขภาพจิตนักเรียน  อนุบาล–มัธยมปลาย  จังหวัดเลย  ประเทศไทย',
    Inches(1.3), Inches(3.15), Inches(11), Inches(0.52),
    size=18, color=LIGHT)
txt(sl, 'มาตรฐานสากล 3 แบบประเมิน  ·  Real-time Dashboard  ·  Crisis Alert อัตโนมัติ  ·  PDPA Compliant  ·  Offline PWA',
    Inches(1.3), Inches(3.72), Inches(11), Inches(0.38),
    size=11, color=MUTED)

# Stats row
sw, sh_ = Inches(2.45), Inches(1.05)
sy_ = Inches(4.7)
stats = [('100,000+','นักเรียนในระบบ'), ('3','แบบประเมินมาตรฐาน'), ('5','ระดับสิทธิ์ RBAC'), ('24/7','Crisis Alert')]
for i, (num, lbl) in enumerate(stats):
    stat(sl, Inches(1.3) + i*(sw + Inches(0.18)), sy_, sw, sh_, num, lbl)

txt(sl, 'Executive & Technical Overview  ·  พ.ศ. 2568',
    Inches(1.3), H - Inches(0.52), Inches(8), Inches(0.34),
    size=9, color=MUTED)

# ═══════════════════════════════════════════════════════════════════════════════
# S02 — ปัญหาที่พบ
# ═══════════════════════════════════════════════════════════════════════════════
sl = slide()
header(sl, 'ปัญหาที่เราพบ — ที่มาของโครงการ', 'ที่มา', n=2)

cw_, ch_ = Inches(5.9), Inches(2.5)
cy1_ = Inches(1.08); cy2_ = Inches(3.72)
cx1_ = M;            cx2_ = M + cw_ + Inches(0.22)

card(sl, cx1_, cy1_, cw_, ch_, RED,
     title='สถานการณ์สุขภาพจิตเด็กและเยาวชน', icon='📊',
     body='• 1 ใน 7 เด็กทั่วโลกมีปัญหาสุขภาพจิต\n• เยาวชนไทยอายุ 10-19 ปี มีความเสี่ยงกว่า 15%\n• ส่วนใหญ่ไม่ได้รับการตรวจคัดกรองหรือช่วยเหลือ\n• การพบล่าช้า → ผลกระทบระยะยาวรุนแรง')
card(sl, cx2_, cy1_, cw_, ch_, YELLOW,
     title='ปัญหาในโรงเรียนปัจจุบัน', icon='🏫',
     body='• ครูแนะแนว 1 คน รับผิดชอบนักเรียนหลายร้อยคน\n• ระบบกระดาษ: ช้า ข้อมูลกระจัดกระจาย ค้นหายาก\n• ไม่มีระบบ Early Warning ระดับจังหวัด\n• รอผลการประเมินแต่ละรอบใช้เวลาหลายวัน')
card(sl, cx1_, cy2_, cw_, ch_, RED,
     title='ความเสี่ยงที่อาจพลาด', icon='⚠️',
     body='• ไม่มีระบบแจ้งเตือนกรณีวิกฤต (suicide risk)\n• ผู้บริหารไม่เห็นภาพรวมระดับอำเภอ-จังหวัด\n• ข้อมูลประวัติ ติดตามแนวโน้มไม่ได้\n• นักเรียนในพื้นที่ห่างไกลถูกมองข้ามง่าย')
card(sl, cx2_, cy2_, cw_, ch_, BLUE,
     title='ข้อกำหนด PDPA ที่ต้องปฏิบัติตาม', icon='🔒',
     body='• ข้อมูลสุขภาพจิต = ข้อมูลอ่อนไหว (Sensitive Data)\n• พ.ร.บ. PDPA บังคับเข้ารหัสและควบคุมการเข้าถึง\n• ระบบกระดาษไม่สามารถรองรับมาตรฐานได้\n• ต้องมี Audit Log ทุกการเข้าถึงข้อมูลนักเรียน')

# ═══════════════════════════════════════════════════════════════════════════════
# S03 — LEMCS คืออะไร
# ═══════════════════════════════════════════════════════════════════════════════
sl = slide()
header(sl, 'LEMCS คืออะไร?', 'ภาพรวมระบบ', n=3)

lw_ = Inches(6.8)
mtxt(sl, [
    {'t': 'LEMCS  (Loei Educational MindCare System)', 's': 18, 'c': TEAL2, 'b': True},
    {'t': ' ', 's': 5},
    {'t': 'คือ แพลตฟอร์มดิจิทัล สำหรับคัดกรองและติดตามสุขภาพจิต', 's': 13, 'c': LIGHT},
    {'t': 'นักเรียนระดับ อนุบาล–มัธยมปลาย ของจังหวัดเลยแบบครบวงจร', 's': 13, 'c': LIGHT},
    {'t': ' ', 's': 5},
    {'t': 'ออกแบบให้ "นักเรียนทำแบบประเมินเองได้" บนมือถือ', 's': 12, 'c': MUTED},
    {'t': 'ขณะที่ครูและผู้บริหารเห็นผลแบบ Real-time', 's': 12, 'c': MUTED},
    {'t': 'พร้อมระบบแจ้งเตือนอัตโนมัติทันทีที่พบความเสี่ยง', 's': 12, 'c': MUTED},
    {'t': ' ', 's': 8},
    {'t': 'คุณสมบัติสำคัญ:', 's': 11, 'c': TEAL2, 'b': True},
    {'t': '✓  Web Application — ไม่ต้องติดตั้งแอปเพิ่มเติม', 's': 11, 'c': LIGHT},
    {'t': '✓  Mobile-first — รองรับมือถือทุกรุ่น ขนาด 375px+', 's': 11, 'c': LIGHT},
    {'t': '✓  Offline PWA — ใช้งานได้แม้ไม่มีอินเตอร์เน็ต', 's': 11, 'c': LIGHT},
    {'t': '✓  PDPA Compliant — เข้ารหัส AES-256 ทุก record', 's': 11, 'c': LIGHT},
    {'t': '✓  Multi-tenant — รองรับขยายไปหลายจังหวัด', 's': 11, 'c': LIGHT},
], M, Inches(1.1), lw_, Inches(5.8))

rx_ = M + lw_ + Inches(0.25)
rw_ = W - rx_ - M
ry_ = Inches(1.1)
for accent_, icon_, title_, body_ in [
    (TEAL,  '🎯', 'กลุ่มเป้าหมาย',
     'นักเรียน อนุบาล–มัธยมปลาย ทั่วจังหวัดเลย\nกว่า 100,000 คน ทุกโรงเรียน ทุกอำเภอ'),
    (BLUE,  '📍', 'ขอบเขตให้บริการ',
     'ครอบคลุมทุกโรงเรียนในจังหวัดเลย\nรองรับขยายไปจังหวัดอื่น (Multi-tenant)'),
    (GREEN, '⚡', 'ช่องทางใช้งาน',
     'เว็บเบราว์เซอร์ บนมือถือ แท็บเล็ต\nหรือคอมพิวเตอร์ — ไม่ต้องติดตั้งแอป'),
    (YELLOW,'🏗️', 'พร้อมใช้งาน',
     'Production URL: lemcs.loeitech.ac.th\nAPI: api-lemcs.loeitech.ac.th'),
]:
    card(sl, rx_, ry_, rw_, Inches(1.45), accent_, title=title_, icon=icon_, body=body_)
    ry_ += Inches(1.55)

# ═══════════════════════════════════════════════════════════════════════════════
# S04 — ผู้ใช้งานระบบ
# ═══════════════════════════════════════════════════════════════════════════════
sl = slide()
header(sl, 'ผู้ใช้งานและระดับสิทธิ์ (RBAC) — 5 ระดับ', 'ผู้ใช้งาน', n=4)

rows_ = [
    ('ผู้ใช้งาน', 'บทบาท', 'สิ่งที่ทำได้ในระบบ', 'ขอบเขตข้อมูล', 'Session'),
    ('🧒 นักเรียน',           'Student',          'ทำแบบประเมิน · ดูผลและคำแนะนำ · ดูประวัติตัวเอง',                        'เฉพาะตนเอง',     '60 นาที'),
    ('👩‍🏫 ครูแนะแนว',          'School Admin',     'ติดตามนักเรียน · รับ Alert · Dashboard โรงเรียน · Export รายงาน',         'ในโรงเรียนตัวเอง','8 ชั่วโมง'),
    ('🏛️ ศึกษานิเทศก์/ผอ.เขต','Commission Admin',  'เปรียบเทียบข้ามโรงเรียน · แนวโน้มระดับอำเภอ · รายงาน',                  'ในอำเภอตัวเอง',  '8 ชั่วโมง'),
    ('🏢 ผู้บริหารจังหวัด',     'Super Admin',      'ภาพรวมทั้งจังหวัด · วิเคราะห์ Trend · Import นักเรียน · จัดการผู้ใช้', 'ทั้งจังหวัด',    '8 ชั่วโมง'),
    ('⚙️ System Admin',        'System Admin',     'ทุกอย่าง + Audit Log + System Config + Multi-tenant mgmt',               'ทุก tenant',     '8 ชั่วโมง'),
]
tbl(sl, M, Inches(1.12), W - 2*M, Inches(6.0), rows_,
    [Inches(2.3), Inches(1.9), Inches(5.3), Inches(2.0), Inches(1.3)])

# ═══════════════════════════════════════════════════════════════════════════════
# S05 — Section: เครื่องมือ
# ═══════════════════════════════════════════════════════════════════════════════
sec_slide('SECTION  01', 'เครื่องมือและฟีเจอร์หลัก',
          'แบบประเมินมาตรฐานสากล  ·  Student Journey  ·  Crisis Alert  ·  Dashboard  ·  Mobile PWA',
          accent=TEAL, n=5)

# ═══════════════════════════════════════════════════════════════════════════════
# S06 — แบบประเมิน 3 ชุด
# ═══════════════════════════════════════════════════════════════════════════════
sl = slide()
header(sl, 'แบบประเมินมาตรฐานสากล 3 ชุด', 'เครื่องมือประเมิน', n=6)

cw3_ = Inches(3.95); ch3_ = Inches(5.45); cy3_ = Inches(1.08); gap3_ = Inches(0.22)
assess_ = [
    (YELLOW, '😰', 'ST-5', 'วัดระดับความเครียด', [
        ('อายุเป้าหมาย', '15 ปีขึ้นไป'),
        ('จำนวนข้อ', '5 ข้อคำถาม'),
        ('ช่วงคะแนน', '0 – 15 คะแนน'),
        ('ระดับผล', '4 ระดับ'),
        ('ปกติ', '≤ 4 คะแนน'),
        ('เล็กน้อย', '5 – 7 คะแนน'),
        ('ปานกลาง', '8 – 9 คะแนน'),
        ('รุนแรง', '≥ 10 คะแนน'),
        ('Suicide Flag', 'ไม่มี'),
        ('มาตรฐาน', 'กรมสุขภาพจิต ไทย'),
    ]),
    (BLUE, '😔', 'PHQ-A', 'ภาวะซึมเศร้าวัยรุ่น', [
        ('อายุเป้าหมาย', '11 – 20 ปี'),
        ('จำนวนข้อ', '9 หลัก + BQ1/BQ2'),
        ('ช่วงคะแนน', '0 – 27 คะแนน'),
        ('ระดับผล', '5 ระดับ'),
        ('None', '0 – 4'),
        ('Mild', '5 – 9'),
        ('Moderate', '10 – 14'),
        ('Severe', '15 – 19'),
        ('Very Severe', '≥ 20'),
        ('Suicide Flag', 'Q9 ≥ 1 หรือ BQ1/BQ2 = true'),
    ]),
    (GREEN, '🧒', 'CDI', 'ภาวะซึมเศร้าเด็ก', [
        ('อายุเป้าหมาย', '7 – 17 ปี'),
        ('จำนวนข้อ', '27 ข้อ (ก/ข/ค)'),
        ('ช่วงคะแนน', '0 – 54 คะแนน'),
        ('ระดับผล', '2 ระดับ'),
        ('ปกติ', '0 – 14 คะแนน'),
        ('มีนัยทางคลินิก', '≥ 15 → แจ้งเตือน'),
        ('Group A', 'ก=0, ข=1, ค=2 (ปกติ)'),
        ('Group B', 'ก=2, ข=1, ค=0 (invert)'),
        ('Suicide Flag', 'ไม่มี — ≥15 triggers alert'),
        ('มาตรฐาน', 'Maria Kovacs 1985, อายุ 7-17'),
    ]),
]

for i, (col_, icon_, name_, desc_, details_) in enumerate(assess_):
    cx3_ = M + i * (cw3_ + gap3_)
    box(sl, cx3_, cy3_, cw3_, ch3_, CARD)
    box(sl, cx3_, cy3_, cw3_, Inches(0.055), col_)
    txt(sl, icon_,  cx3_ + Inches(0.18), cy3_ + Inches(0.12), Inches(0.55), Inches(0.48), size=22)
    txt(sl, name_,  cx3_ + Inches(0.18), cy3_ + Inches(0.65), cw3_ - Inches(0.28), Inches(0.45),
        size=24, color=col_, bold=True)
    txt(sl, desc_,  cx3_ + Inches(0.18), cy3_ + Inches(1.13), cw3_ - Inches(0.28), Inches(0.3),
        size=10, color=MUTED)
    box(sl, cx3_ + Inches(0.18), cy3_ + Inches(1.48), cw3_ - Inches(0.4), Inches(0.02),
        RGBColor(0x1E, 0x34, 0x50))
    dy_ = cy3_ + Inches(1.6)
    for lbl_, val_ in details_:
        txt(sl, f'{lbl_}:', cx3_ + Inches(0.18), dy_, Inches(1.6),  Inches(0.27), size=8,  color=MUTED)
        txt(sl, val_,       cx3_ + Inches(1.72), dy_, cw3_ - Inches(1.85), Inches(0.27), size=8,  color=LIGHT, bold=True)
        dy_ += Inches(0.31)

note_y_ = cy3_ + ch3_ + Inches(0.1)
box(sl, M, note_y_, W - 2*M, Inches(0.35), DKBLUE)
txt(sl, '💡  ระบบเลือกแบบประเมินอัตโนมัติตามอายุ  |  15-17 ปี → ครบ 3 ชุด  |  11-14 ปี → CDI + PHQ-A  |  7-10 ปี → CDI เท่านั้น  |  <7 ปี หรือไม่มีวันเกิด → ไม่แสดง',
    M + Inches(0.18), note_y_ + Inches(0.05), W - 2*M - Inches(0.3), Inches(0.25),
    size=8, color=TEAL)

# ═══════════════════════════════════════════════════════════════════════════════
# S07 — Student Login & Dashboard
# ═══════════════════════════════════════════════════════════════════════════════
sl = slide()
header(sl, 'เส้นทางการใช้งาน — มุมมองนักเรียน: Login + Dashboard', 'การทำงาน', n=7)

# Banner: No OTP
box(sl, M, Inches(1.08), W - 2*M, Inches(0.3), DKBLUE)
txt(sl, 'ขั้นตอนการเข้าสู่ระบบ — ไม่มี OTP ทุกสังกัด  |  เลือกสังกัดก่อน จะได้ฟอร์มที่แตกต่างกัน',
    M + Inches(0.18), Inches(1.1), W - 2*M - Inches(0.3), Inches(0.22),
    size=9, color=TEAL2, bold=True)

# Two login cards
lw7_ = Inches(5.95)
lh7_ = Inches(2.48)
ly7_ = Inches(1.45)

# Left card: สพฐ.
box(sl, M, ly7_, lw7_, lh7_, CARD)
box(sl, M, ly7_, lw7_, Inches(0.05), TEAL)
txt(sl, 'สพฐ. / อาชีวะ / เอกชน', M + Inches(0.15), ly7_ + Inches(0.1),
    lw7_ - Inches(0.2), Inches(0.32), size=13, color=TEAL2, bold=True)
txt(sl, 'POST /api/auth/login/bypass', M + Inches(0.15), ly7_ + Inches(0.44),
    lw7_ - Inches(0.3), Inches(0.2), size=8, color=MUTED)
spb_steps_ = [
    ('1', 'เลือกสังกัด',      'สพฐ. / อาชีวะ / เอกชน'),
    ('2', 'เลขบัตรประชาชน',   '13 หลัก (เข้ารหัส AES-256 ในฐานข้อมูล)'),
    ('3', 'วันเกิด',           'DD/MM/YYYY พ.ศ.'),
    ('4', 'รหัสนักเรียน',      'ออกโดยโรงเรียน (student_code)'),
]
iy7_ = ly7_ + Inches(0.72)
for n7_, fn7_, fd7_ in spb_steps_:
    box(sl, M + Inches(0.14), iy7_, Inches(0.25), Inches(0.25), TEAL)
    txt(sl, n7_, M + Inches(0.14), iy7_, Inches(0.25), Inches(0.25),
        size=8, color=NAVY, bold=True, align=PP_ALIGN.CENTER)
    txt(sl, fn7_, M + Inches(0.46), iy7_ + Inches(0.02), Inches(1.9), Inches(0.2),
        size=10, color=LIGHT, bold=True)
    txt(sl, fd7_, M + Inches(0.46), iy7_ + Inches(0.22), lw7_ - Inches(0.65), Inches(0.18),
        size=8, color=MUTED)
    iy7_ += Inches(0.43)

# Right card: สกร.
rx7_ = M + lw7_ + Inches(0.23)
box(sl, rx7_, ly7_, lw7_, lh7_, CARD)
box(sl, rx7_, ly7_, lw7_, Inches(0.05), BLUE)
txt(sl, 'สกร. (การศึกษานอกระบบ)', rx7_ + Inches(0.15), ly7_ + Inches(0.1),
    lw7_ - Inches(0.2), Inches(0.32), size=13, color=BLUE, bold=True)
txt(sl, 'POST /api/auth/login/skr', rx7_ + Inches(0.15), ly7_ + Inches(0.44),
    lw7_ - Inches(0.3), Inches(0.2), size=8, color=MUTED)
skr_steps_ = [
    ('1', 'เลือกสังกัด',   'สกร.'),
    ('2', 'รหัสนักเรียน', 'ออกโดยสถานศึกษา สกร.'),
    ('3', 'วันเกิด',      'DD/MM/YYYY พ.ศ. (ย้อนหลังได้ 90 ปี)'),
]
iy7_r = ly7_ + Inches(0.72)
for n7_r, fn7_r, fd7_r in skr_steps_:
    box(sl, rx7_ + Inches(0.14), iy7_r, Inches(0.25), Inches(0.25), BLUE)
    txt(sl, n7_r, rx7_ + Inches(0.14), iy7_r, Inches(0.25), Inches(0.25),
        size=8, color=NAVY, bold=True, align=PP_ALIGN.CENTER)
    txt(sl, fn7_r, rx7_ + Inches(0.46), iy7_r + Inches(0.02), Inches(1.9), Inches(0.2),
        size=10, color=LIGHT, bold=True)
    txt(sl, fd7_r, rx7_ + Inches(0.46), iy7_r + Inches(0.22), lw7_ - Inches(0.65), Inches(0.18),
        size=8, color=MUTED)
    iy7_r += Inches(0.43)

# Post-login flow banner
flow_y7_ = ly7_ + lh7_ + Inches(0.12)
box(sl, M, flow_y7_, W - 2*M, Inches(0.28), DKBLUE)
txt(sl, 'หลังเข้าสู่ระบบ — Dashboard Flow',
    M + Inches(0.18), flow_y7_ + Inches(0.04), W - 2*M - Inches(0.3), Inches(0.22),
    size=9, color=TEAL2, bold=True)

# Flow steps
flow_sy7_ = flow_y7_ + Inches(0.32)
flow_sh7_ = H - flow_sy7_ - Inches(0.12)
flow_items7_ = [
    (GREEN,  'PDPA Consent',   'บังคับก่อน\nครั้งแรก\nเท่านั้น'),
    (YELLOW, 'ตรวจรอบสำรวจ', 'Admin เปิด → ทำได้\nปิดรอบ → แสดง\nแจ้งเตือนสีเหลือง'),
    (TEAL,   'Dashboard',      'PendingAssessments\n+ WellnessScore\n+ AssessmentHistory'),
    (BLUE,   'ทำแบบประเมิน',  'ระบบเลือกชุดตามอายุ\nอัตโนมัติ\n(ST-5/PHQ-A/CDI)'),
    (PURPLE, 'ผลลัพธ์ทันที',  'คะแนน + ระดับ\n+ คำแนะนำ\nที่เหมาะสม'),
    (RED,    'Crisis → 1323', 'แสดงสายด่วน\nกรมสุขภาพจิต\nหากพบความเสี่ยง'),
]
n7_ = len(flow_items7_)
fsw7_ = (W - 2*M) / n7_
for i7_, (col7_, t7_, d7_) in enumerate(flow_items7_):
    fx7_ = M + i7_ * fsw7_
    fw7_ = fsw7_ - Inches(0.04)
    box(sl, fx7_, flow_sy7_, fw7_, flow_sh7_, CARD)
    box(sl, fx7_, flow_sy7_, fw7_, Inches(0.05), col7_)
    txt(sl, t7_, fx7_, flow_sy7_ + Inches(0.08), fw7_, Inches(0.32),
        size=9, color=col7_, bold=True, align=PP_ALIGN.CENTER)
    txt(sl, d7_, fx7_, flow_sy7_ + Inches(0.44), fw7_, flow_sh7_ - Inches(0.55),
        size=8, color=MUTED, align=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════════════════════════════
# S08 — Crisis Alert
# ═══════════════════════════════════════════════════════════════════════════════
sl = slide()
header(sl, 'Crisis Alert System — ระบบแจ้งเตือนวิกฤตในแอป', 'ระบบแจ้งเตือน', n=8)

lw8_ = Inches(7.8)

mtxt(sl, [
    {'t': 'ตรวจจับความเสี่ยงทันทีหลังนักเรียนส่งแบบประเมิน', 's': 14, 'c': LIGHT, 'b': True},
    {'t': ' ', 's': 4},
    {'t': 'เงื่อนไขที่ trigger Crisis Alert:', 's': 11, 'c': TEAL2, 'b': True},
    {'t': '• PHQ-A ข้อ 9 (Q9) ≥ 1 คะแนน (คิดเรื่องทำร้ายตัวเอง)', 's': 10, 'c': LIGHT},
    {'t': '• PHQ-A Bonus Q1 หรือ BQ2 = true', 's': 10, 'c': LIGHT},
    {'t': '• CDI คะแนนรวม ≥ 15 (มีนัยสำคัญทางคลินิก)', 's': 10, 'c': LIGHT},
    {'t': ' ', 's': 5},
    {'t': 'กระบวนการแจ้งเตือน (Synchronous — ไม่ผ่านคิว Celery):', 's': 11, 'c': TEAL2, 'b': True},
], M, Inches(1.1), lw8_, Inches(2.2))

steps8_ = [
    (RED,   '1', 'ตรวจพบความเสี่ยง',
     'สร้าง Critical Alert record ทันที\nบันทึก timestamp + ข้อมูลนักเรียน (Synchronous)'),
    (TEAL,  '2', 'In-App Dashboard แจ้งเตือน',
     'ปรากฏใน Admin Dashboard Section B ทันที\nRefresh อัตโนมัติทุก 30 วินาที\nติดตามสถานะ: รอดำเนินการ / กำลังดูแล / เสร็จสิ้น'),
    (GREEN, '3', 'หน้าจอนักเรียน',
     'แสดงสายด่วน 1323 กรมสุขภาพจิต\nพร้อมคำแนะนำขอความช่วยเหลือทันที'),
]

ty8_ = Inches(3.45)
for col8_, num8_, ttl8_, dsc8_ in steps8_:
    box(sl, M, ty8_, Inches(0.36), Inches(0.36), col8_)
    txt(sl, num8_, M, ty8_, Inches(0.36), Inches(0.36),
        size=11, color=NAVY, bold=True, align=PP_ALIGN.CENTER)
    txt(sl, ttl8_, M + Inches(0.46), ty8_ + Inches(0.01), lw8_ - Inches(0.55), Inches(0.26),
        size=11, color=col8_, bold=True)
    txt(sl, dsc8_, M + Inches(0.46), ty8_ + Inches(0.28), lw8_ - Inches(0.55), Inches(0.55),
        size=9, color=MUTED)
    ty8_ += Inches(0.9)

# Right cards
rx8_ = M + lw8_ + Inches(0.28)
rw8_ = W - rx8_ - M
ry8_ = Inches(1.1)
for col8_, icon8_, ttl8_, body8_ in [
    (TEAL, '🔔', 'In-App Dashboard',
     'Admin เห็น Alert ใน Section B\nทันทีที่นักเรียนส่งผล\nRefresh ทุก 30 วินาที'),
    (RED,  '🆘', 'สายด่วน 1323',
     'แสดงบนหน้าจอนักเรียนทันที\nกรมสุขภาพจิต — 24 ชั่วโมง\nพร้อมคำแนะนำ'),
]:
    card(sl, rx8_, ry8_, rw8_, Inches(2.9), col8_, title=ttl8_, icon=icon8_, body=body8_)
    ry8_ += Inches(3.05)

# ═══════════════════════════════════════════════════════════════════════════════
# S09 — Admin Dashboard 4 Sections
# ═══════════════════════════════════════════════════════════════════════════════
sl = slide()
header(sl, 'Admin Dashboard — 4 ส่วนหลักของการวิเคราะห์ข้อมูล', 'Dashboard', n=9)

cw9_ = (W - 2*M - Inches(0.25)) / 2
ch9_ = Inches(2.3)
gap9_ = Inches(0.25)
cy9_1 = Inches(1.08)
cy9_2 = cy9_1 + ch9_ + Inches(0.18)
cx9_1 = M
cx9_2 = M + cw9_ + gap9_

sections9_ = [
    (cx9_1, cy9_1, TEAL,   'A',
     'ส่วน A — ผู้บริหาร (Executive)',
     'StatsCards (refresh 60 วินาที): จำนวนนักเรียน,\nการประเมินวันนี้, ความเสี่ยงสูง, รอดำเนินการ\n'
     'RiskFunnelChart: สัดส่วนความเสี่ยงทุกระดับ\n'
     'InsightPanel: ข้อสรุปอัตโนมัติ\n'
     'AffiliationStudentStats: สถิติแยกสังกัด'),
    (cx9_2, cy9_1, RED,    'B',
     'ส่วน B — ศูนย์ปฏิบัติการ (Operations)',
     'AlertStatusSummary: สรุปสถานะ Alert ทั้งหมด\n'
     'RecentAlerts: 8 รายการล่าสุด (refresh 30 วินาที)\n'
     'Quick Actions: อัปเดตสถานะ Alert ทันที\n'
     'สถานะ: รอดำเนินการ / กำลังดูแล / เสร็จสิ้น'),
    (cx9_1, cy9_2, BLUE,   'C',
     'ส่วน C — วิเคราะห์เชิงลึก (Analytics)',
     'TrendChart: แนวโน้มรายเดือน/สัปดาห์\n'
     'MoMDeltaChart: เปรียบเทียบเดือนต่อเดือน\n'
     'SeverityChart: Doughnut แยกระดับความเสี่ยง\n'
     'AssessmentTypeChart: Stacked Bar แยกประเภท\n'
     'RiskProgressBars: Progress bar รายระดับ'),
    (cx9_2, cy9_2, PURPLE, 'D',
     'ส่วน D — เปรียบเทียบองค์กร (Comparison)',
     'OrgCompareChart: 100% Stacked Bar\n'
     'Auto group: district → schools,\n'
     'affiliation → districts, else → affiliations\n'
     'ดาวน์โหลด PNG ได้ทุก chart (DownloadBtn)'),
]

for cx9_, cy9_, col9_, lbl9_, ttl9_, dsc9_ in sections9_:
    box(sl, cx9_, cy9_, cw9_, ch9_, CARD)
    box(sl, cx9_, cy9_, cw9_, Inches(0.055), col9_)
    box(sl, cx9_, cy9_, Inches(0.38), Inches(0.38), col9_)
    txt(sl, lbl9_, cx9_, cy9_, Inches(0.38), Inches(0.38),
        size=14, color=NAVY, bold=True, align=PP_ALIGN.CENTER)
    txt(sl, ttl9_, cx9_ + Inches(0.46), cy9_ + Inches(0.08),
        cw9_ - Inches(0.55), Inches(0.3), size=11, color=col9_, bold=True)
    mtxt(sl, [{'t': line, 's': 8, 'c': MUTED} for line in dsc9_.split('\n')],
         cx9_ + Inches(0.18), cy9_ + Inches(0.44), cw9_ - Inches(0.28), ch9_ - Inches(0.52))

# Filter + Export footer
fy9_ = cy9_2 + ch9_ + Inches(0.16)
fh9_ = H - fy9_ - Inches(0.1)
box(sl, M, fy9_, W - 2*M, fh9_, DKBLUE)
txt(sl, 'FilterBar:', M + Inches(0.18), fy9_ + Inches(0.1),
    Inches(1.1), Inches(0.25), size=9, color=TEAL2, bold=True)
txt(sl, 'survey_round_id  |  affiliation_id  |  district_id  |  school_id  |  assessment_type  |  grade  |  gender  |  date_from  |  date_to',
    M + Inches(1.3), fy9_ + Inches(0.1), W - 2*M - Inches(1.5), Inches(0.25),
    size=8, color=LIGHT)
txt(sl, 'Export:', M + Inches(0.18), fy9_ + Inches(0.4),
    Inches(0.9), Inches(0.25), size=9, color=TEAL2, bold=True)
txt(sl, 'PDF (lemcs_report.pdf)  |  Excel (lemcs_report.xlsx)  |  PNG per chart (DownloadBtn)  |  Print',
    M + Inches(1.3), fy9_ + Inches(0.4), W - 2*M - Inches(1.5), Inches(0.25),
    size=8, color=LIGHT)

# ═══════════════════════════════════════════════════════════════════════════════
# S10 — Mobile PWA
# ═══════════════════════════════════════════════════════════════════════════════
sl = slide()
header(sl, 'Mobile PWA — ใช้งานได้ทุกที่ แม้ไม่มีอินเตอร์เน็ต', 'Mobile PWA', n=10)

lw10_ = Inches(7.5)
mtxt(sl, [
    {'t': 'Progressive Web App (PWA)', 's': 17, 'c': TEAL2, 'b': True},
    {'t': ' ', 's': 4},
    {'t': 'ไม่ต้องดาวน์โหลดจาก App Store — เปิดเว็บแล้วใช้ได้ทันที', 's': 12, 'c': LIGHT},
    {'t': 'เพิ่มไอคอนหน้า Home Screen ได้เหมือนแอปทั่วไป', 's': 12, 'c': LIGHT},
    {'t': ' ', 's': 5},
    {'t': 'ฟีเจอร์ Offline:', 's': 11, 'c': TEAL2, 'b': True},
    {'t': '• Service Worker cache คำถาม 3 แบบประเมินไว้ในเครื่อง', 's': 11, 'c': LIGHT},
    {'t': '• ทำแบบประเมินได้แม้สัญญาณขาด หรืออยู่ในโรงเรียนห่างไกล', 's': 11, 'c': LIGHT},
    {'t': '• ระบบ Sync ข้อมูลอัตโนมัติเมื่อกลับออนไลน์', 's': 11, 'c': LIGHT},
    {'t': ' ', 's': 5},
    {'t': 'การรองรับอุปกรณ์:', 's': 11, 'c': TEAL2, 'b': True},
    {'t': '• Mobile-first design — รองรับหน้าจอ 375px ขึ้นไป', 's': 11, 'c': LIGHT},
    {'t': '• ทดสอบบน Android, iOS (Chrome, Safari)', 's': 11, 'c': LIGHT},
    {'t': '• Font ภาษาไทย: Noto Sans Thai / Sarabun', 's': 11, 'c': LIGHT},
], M, Inches(1.1), lw10_, Inches(5.0))

rx10_ = M + lw10_ + Inches(0.22)
rw10_ = W - rx10_ - M
ry10_ = Inches(1.1)
for col10_, icon10_, lbl10_, dsc10_ in [
    (TEAL,  '📱', 'Mobile-first',  'หน้าจอ 375px+ ใช้งานง่าย\nบนมือถือทุกรุ่น'),
    (GREEN, '✈️', 'Offline Mode',  'ทำแบบประเมินได้\nแม้ไม่มีอินเตอร์เน็ต'),
    (BLUE,  '⚡', 'Fast Load',    'Cached assets โหลดเร็ว\nแม้เน็ตช้า'),
    (YELLOW,'🔄', 'Auto Sync',    'Sync อัตโนมัติ\nเมื่อกลับออนไลน์'),
]:
    card(sl, rx10_, ry10_, rw10_, Inches(1.42), col10_,
         title=lbl10_, icon=icon10_, body=dsc10_)
    ry10_ += Inches(1.52)

# ═══════════════════════════════════════════════════════════════════════════════
# S11 — Section: ความปลอดภัย
# ═══════════════════════════════════════════════════════════════════════════════
sec_slide('SECTION  02', 'ความปลอดภัยและ PDPA',
          'AES-256 Encryption  ·  RBAC 5 ระดับ  ·  Audit Log  ·  JWT Session',
          accent=GREEN, n=11)

# ═══════════════════════════════════════════════════════════════════════════════
# S12 — PDPA & Security
# ═══════════════════════════════════════════════════════════════════════════════
sl = slide()
header(sl, 'ความปลอดภัยข้อมูลและ PDPA Compliance', 'ความปลอดภัย', n=12)

sec12_ = [
    (TEAL,   '🔐', 'AES-256 Encryption',
     'เลขบัตรประชาชนเข้ารหัส AES-256 ทุก record\nไม่ปรากฏ plain text ใน DB, URL, หรือ Log\nEncryption Key ต้องยาว 32 ตัวอักษร (256-bit)\nระบบ crash ทันทีหาก Key ไม่ครบ 32 ตัว'),
    (GREEN,  '✅', 'PDPA Consent',
     'หน้า Consent บังคับก่อนทำแบบประเมินครั้งแรก\nบันทึกการยินยอม + Timestamp + IP address\nนักเรียนขอถอนความยินยอมได้ตามสิทธิ์\nไม่มี PII ในรูปภาพ URL หรือ query string'),
    (BLUE,   '📋', 'Audit Log',
     'บันทึกทุกการเข้าถึงข้อมูลนักเรียน\nเก็บ user_id + timestamp + action + resource\nรองรับการตรวจสอบตาม PDPA\nค้นหาและ Export ได้โดย System Admin'),
    (YELLOW, '🛡️', 'Session & Auth',
     'นักเรียน: JWT หมดอายุ 60 นาที\nAdmin ทุกระดับ: JWT หมดอายุ 8 ชั่วโมง\nLogin: bypass (สพฐ./อาชีวะ/เอกชน) หรือ skr — ไม่มี OTP\nRate limit 100 req/นาที ต่อ IP ผ่าน Redis'),
    (PURPLE, '🚫', 'No PII in URL',
     'ไม่มีชื่อ เลขบัตร หรือข้อมูลส่วนตัวใน URL\nทุก endpoint ใช้ UUID ไม่ใช่ชื่อหรือ ID จริง\nป้องกันการ expose ผ่าน Browser history\nหรือ Server access log'),
    (RED,    '🔄', 'Backup & Recovery',
     'pg_dump อัตโนมัติทุกวันเวลา 02:00\nส่งไป MinIO Object Storage (encrypted)\nรองรับ Point-in-time recovery\nตรวจสอบ Backup integrity อัตโนมัติ'),
]

sw12_ = Inches(3.95); sh12_ = Inches(2.5)
sy12_1 = Inches(1.08); sy12_2 = Inches(3.72); sg12_ = Inches(0.22)

for i12_, (col12_, icon12_, ttl12_, dsc12_) in enumerate(sec12_):
    sx12_ = M + (i12_ % 3) * (sw12_ + sg12_)
    sy12_ = sy12_1 if i12_ < 3 else sy12_2
    card(sl, sx12_, sy12_, sw12_, sh12_, col12_, title=ttl12_, icon=icon12_, body=dsc12_)

# ═══════════════════════════════════════════════════════════════════════════════
# S13 — RBAC
# ═══════════════════════════════════════════════════════════════════════════════
sl = slide()
header(sl, 'Role-Based Access Control (RBAC) — ระดับสิทธิ์และขอบเขตข้อมูล', 'ความปลอดภัย', n=13)

rbac_ = [
    ('บทบาท',                    'Scope ข้อมูล',    'อ่านได้',                                  'เขียน/แก้ไข',                   'ไม่สามารถ'),
    ('🧒 Student',               'ตัวเอง',          'ผลตัวเอง, ประวัติ',                         'ทำแบบประเมิน',                  'ดูคนอื่น'),
    ('👩‍🏫 School Admin',          'โรงเรียนตัวเอง',  'นักเรียนทุกคนในโรงเรียน, Alerts',           'ดูแล Alert, Export รายงาน',     'ข้ามโรงเรียน'),
    ('🏛️ Commission Admin',       'ระดับอำเภอ',      'ทุกโรงเรียนในอำเภอ, รายงานรวม',             'Export รายงาน, ดู Trend',       'ข้ามอำเภอ'),
    ('🏢 Super Admin',            'ทั้งจังหวัด',     'ทุกโรงเรียน, ทุก Alert, ทุกรายงาน',         'จัดการผู้ใช้, Import CSV',      'ข้าม tenant'),
    ('⚙️ System Admin',           'ทุก tenant',      'ทุกอย่าง + Audit Log + System Config',       'จัดการ tenant, ตั้งค่าระบบ',   'ไม่มี'),
]
tbl(sl, M, Inches(1.12), W - 2*M, Inches(6.05), rbac_,
    [Inches(2.1), Inches(1.6), Inches(3.5), Inches(2.8), Inches(2.0)])

# ═══════════════════════════════════════════════════════════════════════════════
# S14 — Section: สถาปัตยกรรม
# ═══════════════════════════════════════════════════════════════════════════════
sec_slide('SECTION  03', 'สถาปัตยกรรมและเทคโนโลยี',
          'Next.js 14  ·  FastAPI Python 3.12  ·  PostgreSQL 16  ·  Redis 7  ·  MinIO  ·  Docker',
          accent=BLUE, n=14)

# ═══════════════════════════════════════════════════════════════════════════════
# S15 — Tech Stack
# ═══════════════════════════════════════════════════════════════════════════════
sl = slide()
header(sl, 'สถาปัตยกรรมระบบ (System Architecture)', 'เทคโนโลยี', n=15)

layers15_ = [
    (TEAL,   '🖥️  Frontend (ผู้ใช้งานโต้ตอบ)',
     'Next.js 14 App Router  ·  TypeScript  ·  Tailwind CSS + DaisyUI  ·  Chart.js (กราฟ)  ·  SWR (data fetching)  ·  next-pwa (Offline)'),
    (BLUE,   '⚙️  Backend (ประมวลผล / API)',
     'FastAPI Python 3.12  ·  SQLAlchemy Async (ORM)  ·  Alembic (DB migration)  ·  Celery (background tasks)  ·  JWT + OTP (Auth)  ·  Axios client'),
    (GREEN,  '🗄️  Data Layer (จัดเก็บข้อมูล)',
     'PostgreSQL 16 (ฐานข้อมูลหลัก)  ·  PgBouncer (connection pool)  ·  Redis 7 (cache + sessions)  ·  MinIO (object storage, backup)'),
    (PURPLE, '🐳  Infrastructure (ระบบโครงสร้าง)',
     'Docker Compose  ·  Nginx Proxy Manager (reverse proxy + SSL)  ·  Grafana (monitoring)  ·  Portainer (container mgmt)  ·  Daily Backup Script'),
]

ly15_ = Inches(1.12)
for col15_, ttl15_, dsc15_ in layers15_:
    box(sl, M, ly15_, W - 2*M, Inches(1.28), CARD)
    box(sl, M, ly15_, Inches(0.055), Inches(1.28), col15_)
    txt(sl, ttl15_, M + Inches(0.18), ly15_ + Inches(0.1), Inches(4.2), Inches(0.35),
        size=12, color=col15_, bold=True)
    txt(sl, dsc15_, M + Inches(0.18), ly15_ + Inches(0.5), W - 2*M - Inches(0.28), Inches(0.55),
        size=10, color=MUTED)
    ly15_ += Inches(1.38)

box(sl, M, ly15_, W - 2*M, Inches(0.55), DKBLUE)
txt(sl, '🌐  Production:  Frontend → lemcs.loeitech.ac.th  ·  API → api-lemcs.loeitech.ac.th  ·  Grafana → grafana-lemcs.loeitech.ac.th  ·  MinIO → minio-lemcs.loeitech.ac.th',
    M + Inches(0.18), ly15_ + Inches(0.14), W - 2*M - Inches(0.28), Inches(0.3),
    size=9, color=TEAL)

# ═══════════════════════════════════════════════════════════════════════════════
# S16 — Database Schema
# ═══════════════════════════════════════════════════════════════════════════════
sl = slide()
header(sl, 'โครงสร้างฐานข้อมูล — Multi-tenant, 9+2 Tables', 'สถาปัตยกรรม', n=16)

lw16_ = Inches(5.8)
mtxt(sl, [
    {'t': 'Multi-tenant Architecture', 's': 15, 'c': TEAL2, 'b': True},
    {'t': ' ', 's': 4},
    {'t': '• system.* — tenants (registry จังหวัด) + super_admins', 's': 11, 'c': LIGHT},
    {'t': '• loei.* — ข้อมูลทั้งหมดของจังหวัดเลย (9 tables)', 's': 11, 'c': LIGHT},
    {'t': '• ขยายรองรับจังหวัดอื่นได้โดยเพิ่ม schema ใหม่', 's': 11, 'c': LIGHT},
    {'t': ' ', 's': 6},
    {'t': 'Partitioning', 's': 13, 'c': TEAL2, 'b': True},
    {'t': '• ตาราง assessments แบ่ง Range Partition ตามปี', 's': 11, 'c': LIGHT},
    {'t': '• Query ข้อมูลปีเดียวเร็วกว่า Full Table Scan มาก', 's': 11, 'c': LIGHT},
    {'t': '• ต้องเพิ่ม Partition ใหม่ทุกปี (Alembic migration)', 's': 11, 'c': LIGHT},
    {'t': ' ', 's': 6},
    {'t': 'Connection Pool', 's': 13, 'c': TEAL2, 'b': True},
    {'t': '• PgBouncer ในโหมด Transaction', 's': 11, 'c': LIGHT},
    {'t': '• รองรับ concurrent users จำนวนมากโดยไม่ล้น DB pool', 's': 11, 'c': LIGHT},
], M, Inches(1.1), lw16_, Inches(5.6))

rx16_ = M + lw16_ + Inches(0.28)
rw16_ = W - rx16_ - M
schema16_ = [
    (TEAL,   'affiliations',  'สังกัดการศึกษา (จังหวัด/สพท.)'),
    (BLUE,   'districts',     'อำเภอ — เชื่อมกับ affiliations'),
    (GREEN,  'schools',       'โรงเรียน — เชื่อมกับ districts'),
    (YELLOW, 'students',      'นักเรียน — National ID (AES-256)'),
    (YELLOW, 'users',         'Admin/Teacher accounts'),
    (PURPLE, 'assessments',   'ผลประเมิน — partitioned by year'),
    (RED,    'alerts',        'การแจ้งเตือนความเสี่ยง'),
    (MUTED,  'notifications', 'ประวัติการส่งแจ้งเตือน'),
    (TEAL,   'audit_logs',    'PDPA Audit — ทุกการเข้าถึง'),
]

ty16_ = Inches(1.1)
for col16_, tname16_, tdesc16_ in schema16_:
    box(sl, rx16_, ty16_, rw16_, Inches(0.5), CARD)
    box(sl, rx16_, ty16_, Inches(0.05), Inches(0.5), col16_)
    txt(sl, tname16_, rx16_ + Inches(0.15), ty16_ + Inches(0.04),
        Inches(2.2), Inches(0.22), size=10, color=col16_, bold=True)
    txt(sl, tdesc16_, rx16_ + Inches(0.15), ty16_ + Inches(0.26),
        rw16_ - Inches(0.2), Inches(0.2), size=8, color=MUTED)
    ty16_ += Inches(0.57)

# ═══════════════════════════════════════════════════════════════════════════════
# S17 — Section: แผนดำเนินการ
# ═══════════════════════════════════════════════════════════════════════════════
sec_slide('SECTION  04', 'แผนดำเนินการและ KPIs',
          'Roadmap  ·  สิ่งที่ต้องการ  ·  ตัวชี้วัดความสำเร็จ  ·  ประโยชน์ที่ได้รับ',
          accent=YELLOW, n=17)

# ═══════════════════════════════════════════════════════════════════════════════
# S18 — Roadmap
# ═══════════════════════════════════════════════════════════════════════════════
sl = slide()
header(sl, 'Roadmap — แผนดำเนินการ 4 Phase', 'แผนงาน', n=18)

phases18_ = [
    (GREEN,  '✓ เสร็จแล้ว', 'Phase 1-3\nระบบหลัก',
     ['Backend API ครบทุก endpoint',
      'ฐานข้อมูล + Multi-tenant Schema',
      'Authentication: OTP + JWT',
      'แบบประเมิน 3 ชุด (ST-5/PHQ-A/CDI)',
      'Admin Dashboard ทุกหน้า',
      'Import ข้อมูลนักเรียน (CSV)',
      'Frontend ทุกหน้า + Offline PWA']),
    (TEAL,   '⏳ กำลังพัฒนา', 'Phase 4\nAlerts & Reports',
     ['In-App Alert system',
      'รายงาน PDF/Excel export',
      'Alert Dashboard + Status tracking',
      'FilterBar & Advanced search',
      'Recommendation engine',
      'System integration test',
      'UAT กับผู้ใช้งานจริง']),
    (BLUE,   '📋 ถัดไป', 'Phase 5\nPilot & Roll-out',
     ['ทดสอบโรงเรียนนำร่อง 3-5 แห่ง',
      'เก็บ Feedback จากครูแนะแนว',
      'ปรับปรุงระบบตาม feedback',
      'อบรมครูแนะแนวทั่วจังหวัด',
      'Import นักเรียนทั้งจังหวัด',
      'เปิดใช้งานจริง (Go-live)',
      'Monitor 30 วันแรก']),
    (YELLOW, '🔮 อนาคต', 'Phase 6\nExpansion',
     ['Multi-province (จังหวัดอื่น)',
      'AI Recommendations',
      'Telehealth integration',
      'Mobile App native',
      'API เชื่อมระบบสาธารณสุข',
      'BI Dashboard ระดับประเทศ',
      'Predictive Analytics']),
]

pw18_ = Inches(2.9); ph18_ = Inches(5.5); py18_ = Inches(1.08); pg18_ = Inches(0.18)

for i18_, (col18_, badge18_, ttl18_, items18_) in enumerate(phases18_):
    px18_ = M + i18_ * (pw18_ + pg18_)
    box(sl, px18_, py18_, pw18_, ph18_, CARD)
    box(sl, px18_, py18_, pw18_, Inches(0.05), col18_)
    box(sl, px18_ + Inches(0.12), py18_ + Inches(0.12), pw18_ - Inches(0.24), Inches(0.3), DKBLUE)
    txt(sl, badge18_, px18_ + Inches(0.12), py18_ + Inches(0.12), pw18_ - Inches(0.24), Inches(0.3),
        size=8, color=col18_, bold=True, align=PP_ALIGN.CENTER)
    for line18_ in ttl18_.split('\n'):
        txt(sl, line18_, px18_ + Inches(0.15), py18_ + Inches(0.55 if line18_ == ttl18_.split('\n')[0] else 0.85),
            pw18_ - Inches(0.25), Inches(0.35),
            size=12 if '\n' not in line18_ else 10,
            color=WHITE if '\n' not in line18_ else col18_, bold=True)
    iy18_ = py18_ + Inches(1.3)
    for item18_ in items18_:
        txt(sl, f'• {item18_}', px18_ + Inches(0.15), iy18_, pw18_ - Inches(0.25), Inches(0.35),
            size=9, color=MUTED)
        iy18_ += Inches(0.44)

# ═══════════════════════════════════════════════════════════════════════════════
# S19 — สิ่งที่ต้องการ
# ═══════════════════════════════════════════════════════════════════════════════
sl = slide()
header(sl, 'สิ่งที่ต้องการเพื่อเริ่มใช้งานจริง', 'ความต้องการ', n=19)

needs19_ = [
    (TEAL,   '📁', 'ข้อมูลนักเรียน CSV',
     'รายชื่อนักเรียนทุกโรงเรียน\n• รหัสนักเรียน, เลขบัตรประชาชน\n• วันเกิด, ชั้นเรียน, โรงเรียน\nระบบรองรับ Import CSV จาก Excel ได้เลย'),
    (GREEN,  '👤', 'บัญชีผู้ใช้งาน Admin',
     'สร้างบัญชีสำหรับผู้ดูแลระบบแต่ละระดับ\n• School Admin: ครูแนะแนวแต่ละโรงเรียน\n• Commission Admin: ศึกษานิเทศก์/ผอ.เขต\n• Super Admin: ผู้บริหารจังหวัด'),
    (BLUE,   '📅', 'เปิดรอบสำรวจ (Survey Round)',
     'Admin ต้องสร้าง Survey Round ก่อน\nนักเรียนจึงจะเห็นและทำแบบประเมินได้\n• กำหนดปีการศึกษา + ภาคเรียน\n• เปิด/ปิด ได้ตลอดเวลา'),
    (YELLOW, '👥', 'ทีมทดสอบนำร่อง',
     'ผู้ร่วมทดสอบช่วง Pilot 2-4 สัปดาห์\n• ครูแนะแนว 1-2 คน/โรงเรียน (3-5 โรง)\n• นักเรียกลุ่มทดสอบ 20-50 คน\n• ผู้บริหาร UAT และให้ Feedback'),
    (PURPLE, '🔑', 'Credentials & Server',
     'ข้อมูลสำหรับตั้งค่าและ Deploy\n• Encryption Key ยาว 32 ตัวอักษร (AES-256)\n• Server หรือ Cloud (Docker Compose ready)\n• Domain + SSL Certificate'),
    (RED,    '📅', 'แผนอบรมผู้ใช้งาน',
     'เตรียมผู้ใช้งานก่อน Roll-out\n• อบรมครูแนะแนว: Dashboard + Alert\n• อบรมผู้บริหาร: รายงานและ Trend\n• เอกสาร User Manual ภาษาไทย'),
]

nw19_ = Inches(3.95); nh19_ = Inches(2.55); ng19_ = Inches(0.22)
ny19_1 = Inches(1.1); ny19_2 = Inches(3.78)

for i19_, (col19_, icon19_, ttl19_, dsc19_) in enumerate(needs19_):
    nx19_ = M + (i19_ % 3) * (nw19_ + ng19_)
    ny19_ = ny19_1 if i19_ < 3 else ny19_2
    card(sl, nx19_, ny19_, nw19_, nh19_, col19_, title=ttl19_, icon=icon19_, body=dsc19_)

# ═══════════════════════════════════════════════════════════════════════════════
# S20 — KPIs
# ═══════════════════════════════════════════════════════════════════════════════
sl = slide()
header(sl, 'ตัวชี้วัดความสำเร็จ (Key Performance Indicators)', 'KPIs', n=20)

# Left KPIs
txt(sl, '📊  KPI ด้านการใช้งาน (Operational)', M, Inches(1.12), Inches(6), Inches(0.38),
    size=14, color=TEAL2, bold=True)

op20_ = [
    (TEAL,   'ความครอบคลุม',    '≥ 80%',    'นักเรียนทำแบบประเมินครบภายใน 30 วันแรก'),
    (BLUE,   'Alert Response',  '< 5 นาที', 'เวลาแจ้งเตือนถึงครูนับจากที่นักเรียนส่งผล'),
    (GREEN,  'System Uptime',   '≥ 99.5%',  'ระบบพร้อมใช้งานตลอดเวลาทำการโรงเรียน'),
    (YELLOW, 'Alert Follow-up', '100%',     'ทุก Critical Alert ต้องมีการดำเนินการ'),
    (PURPLE, 'Data Accuracy',   '100%',     'ผลการประเมินตรงกับ scoring algorithm สากล'),
]

ky20_ = Inches(1.6)
for col20_, kname20_, target20_, kdesc20_ in op20_:
    box(sl, M, ky20_, Inches(6.0), Inches(0.78), CARD)
    box(sl, M, ky20_, Inches(0.05), Inches(0.78), col20_)
    txt(sl, kname20_,  M + Inches(0.18), ky20_ + Inches(0.07), Inches(2.0), Inches(0.28), size=11, color=col20_, bold=True)
    txt(sl, target20_, M + Inches(2.3),  ky20_ + Inches(0.03), Inches(1.4), Inches(0.38), size=20, color=col20_, bold=True)
    txt(sl, kdesc20_,  M + Inches(0.18), ky20_ + Inches(0.42), Inches(5.7), Inches(0.28), size=9, color=MUTED)
    ky20_ += Inches(0.87)

# Right KPIs
rx20_ = M + Inches(6.4)
rw20_ = W - rx20_ - M
txt(sl, '🎯  KPI ด้านผลลัพธ์ (Outcome)', rx20_, Inches(1.12), rw20_, Inches(0.38),
    size=14, color=TEAL2, bold=True)

out20_ = [
    (TEAL,   'Early Detection',   'พบเคสวิกฤตก่อนบานปลาย — ลดผลกระทบระยะยาว'),
    (GREEN,  'ลดเวลาประเมิน',    'จากหลายวัน → ทันที — ครูมีเวลาดูแลมากขึ้น'),
    (BLUE,   'Data-driven Policy','ผู้บริหารตัดสินใจจากข้อมูลจริง ไม่ใช่ความรู้สึก'),
    (YELLOW, 'PDPA Compliant',    'ผ่านการตรวจสอบ 100% ลดความเสี่ยงทางกฎหมาย'),
    (RED,    'ครอบคลุมห่างไกล',  'โรงเรียนพื้นที่ห่างไกล Offline PWA ยังใช้ได้'),
]

ky20_r = Inches(1.6)
for col20_r, kname20_r, kdesc20_r in out20_:
    box(sl, rx20_, ky20_r, rw20_, Inches(0.78), CARD)
    box(sl, rx20_, ky20_r, Inches(0.05), Inches(0.78), col20_r)
    txt(sl, kname20_r, rx20_ + Inches(0.18), ky20_r + Inches(0.07), rw20_ - Inches(0.25), Inches(0.28),
        size=11, color=col20_r, bold=True)
    txt(sl, kdesc20_r, rx20_ + Inches(0.18), ky20_r + Inches(0.42), rw20_ - Inches(0.25), Inches(0.28),
        size=9, color=MUTED)
    ky20_r += Inches(0.87)

# ═══════════════════════════════════════════════════════════════════════════════
# S21 — ประโยชน์
# ═══════════════════════════════════════════════════════════════════════════════
sl = slide()
header(sl, 'ประโยชน์ที่จังหวัดเลยได้รับ', 'ประโยชน์', n=21)

bens21_ = [
    (TEAL,   '⚡', 'ลดเวลาประเมิน',
     'ประเมิน 100K นักเรียนจากหลายเดือน → ไม่กี่สัปดาห์\nผลรู้ทันที ไม่ต้องรอรวม Excel หรือนับมือ'),
    (GREEN,  '🔭', 'เห็นภาพจังหวัด',
     'ผู้บริหารเห็น Trend รายอำเภอ/โรงเรียน Real-time\nเป็นครั้งแรกของจังหวัด — ตัดสินใจได้ทันที'),
    (BLUE,   '🛡️', 'ไม่พลาดวิกฤต',
     'แจ้งเตือนอัตโนมัติทันที ไม่มีช่องว่างในการดูแล\nลดความเสี่ยงที่จะพลาดเคสที่ต้องการความช่วยเหลือ'),
    (YELLOW, '📁', 'ถูกกฎหมาย PDPA',
     'จัดเก็บและประมวลผลข้อมูลถูกต้องตามกฎหมาย\nลดความเสี่ยงทางกฎหมาย เบี้ยปรับสูงสุด 5 ล้านบาท'),
    (PURPLE, '📈', 'ข้อมูลเชิงนโยบาย',
     'รายงานประจำปีและ Trend ช่วยวางแผนงบประมาณ\nอ้างอิงในรายงานระดับจังหวัดและส่วนกลาง'),
    (RED,    '📡', 'ครอบคลุมทุกพื้นที่',
     'โรงเรียนห่างไกลไม่มีเน็ตก็ทำแบบประเมินได้\nOffline PWA — Sync อัตโนมัติเมื่อกลับออนไลน์'),
]

bw21_ = Inches(3.95); bh21_ = Inches(2.5); bg21_ = Inches(0.22)
by21_1 = Inches(1.1); by21_2 = Inches(3.72)

for i21_, (col21_, icon21_, ttl21_, dsc21_) in enumerate(bens21_):
    bx21_ = M + (i21_ % 3) * (bw21_ + bg21_)
    by21_ = by21_1 if i21_ < 3 else by21_2
    card(sl, bx21_, by21_, bw21_, bh21_, col21_, title=ttl21_, icon=icon21_, body=dsc21_)

# ═══════════════════════════════════════════════════════════════════════════════
# S22 — Thank You
# ═══════════════════════════════════════════════════════════════════════════════
sl = slide(NAVY2)
box(sl, 0,            0, Inches(0.55), H, TEAL)
box(sl, W - Inches(0.55), 0, Inches(0.55), H, TEAL)

txt(sl, 'ขอบคุณ', W/2 - Inches(1.5), Inches(1.5), Inches(3), Inches(0.55),
    size=16, color=MUTED, align=PP_ALIGN.CENTER)
txt(sl, 'LEMCS', W/2 - Inches(4.5), Inches(1.98), Inches(9), Inches(1.75),
    size=90, color=TEAL2, bold=True, align=PP_ALIGN.CENTER)
txt(sl, 'ระบบดิจิทัลประเมินสุขภาพจิตนักเรียน  อนุบาล–มัธยมปลาย  จังหวัดเลย  ประเทศไทย',
    W/2 - Inches(5.5), Inches(3.65), Inches(11), Inches(0.5),
    size=15, color=LIGHT, align=PP_ALIGN.CENTER)

box(sl, W/2 - Inches(5.2), Inches(4.5), Inches(4.5), Inches(0.88), CARD)
box(sl, W/2 - Inches(5.2), Inches(4.5), Inches(4.5), Inches(0.05), TEAL)
txt(sl, '🌐  เว็บไซต์',         W/2 - Inches(5.0), Inches(4.58), Inches(4.1), Inches(0.26), size=9, color=MUTED)
txt(sl, 'lemcs.loeitech.ac.th', W/2 - Inches(5.0), Inches(4.85), Inches(4.1), Inches(0.4), size=13, color=TEAL2, bold=True)

box(sl, W/2 + Inches(0.7), Inches(4.5), Inches(4.5), Inches(0.88), CARD)
box(sl, W/2 + Inches(0.7), Inches(4.5), Inches(4.5), Inches(0.05), BLUE)
txt(sl, '🔧  API Documentation',          W/2 + Inches(0.9), Inches(4.58), Inches(4.1), Inches(0.26), size=9, color=MUTED)
txt(sl, 'api-lemcs.loeitech.ac.th/docs',  W/2 + Inches(0.9), Inches(4.85), Inches(4.1), Inches(0.4), size=12, color=BLUE, bold=True)

txt(sl, 'พัฒนาเพื่อนักเรียนจังหวัดเลยกว่า 100,000 คน  ·  พ.ศ. 2568',
    W/2 - Inches(5.5), H - Inches(0.58), Inches(11), Inches(0.35),
    size=9, color=MUTED, align=PP_ALIGN.CENTER)

# ── Save ──────────────────────────────────────────────────────────────────────
out = r'D:\@LEMCS\lemcs-presentation.pptx'
prs.save(out)
print(f'[OK] Saved -> {out}')
print(f'   Total slides: {len(prs.slides)}')
